import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import json
from pptx.util import Inches
import matplotlib.pyplot as plt
import io


def read_config(filename):
    try:
        with open(filename) as file:
            return json.load(file)
    except FileNotFoundError:
        return {}

def add_title_slide(pres, conf):
    """
    Creates title slide in pres based on conf
    
    Prameters: 
    pres (Presentation): pptx Presentation object
    conf (dict): parameters for slide creation
    """
    title_slide_layout = pres.slide_layouts[0]
    slide = pres.slides.add_slide(title_slide_layout)
    t = slide.shapes.title
    t.text = conf["title"]
    content = slide.placeholders[1]
    content.text = conf["content"]
    return pres

def add_list_slide(pres, conf):
    """
    Creates list slide in pres based on conf
    
    Prameters: 
    pres (Presentation): pptx Presentation object
    conf (dict): parameters for slide creation
    """
    title_slide_layout = pres.slide_layouts[1]
    slide = pres.slides.add_slide(title_slide_layout)
    t = slide.shapes.title
    t.text = conf["title"]
    content = slide.placeholders[1]
    tf = content.text_frame
    for item in conf["content"]:
        p = tf.add_paragraph()
        p.text = item["text"]
        p.level = int(item["level"])
    return pres

def add_text_slide(pres, conf):
    """
    Creates text slide in pres based on conf
    
    Prameters: 
    pres (Presentation): pptx Presentation object
    conf (dict): parameters for slide creation
    """
    title_slide_layout = pres.slide_layouts[1]
    slide = pres.slides.add_slide(title_slide_layout)
    t = slide.shapes.title
    t.text = conf["title"]
    content = slide.placeholders[1]
    content.text = conf["content"]
    return pres

def add_picture_slide(pres, conf):
    """
    Creates picture slide in pres based on conf
    
    Prameters: 
    pres (Presentation): pptx Presentation object
    conf (dict): parameters for slide creation
    """
    title_slide_layout = pres.slide_layouts[1]
    slide = pres.slides.add_slide(title_slide_layout)
    t = slide.shapes.title
    t.text = conf["title"]
    content = slide.shapes
    picture = content.add_picture(conf["content"], Inches(0.5), Inches(1.75), width=Inches(9), height=Inches(5))
    return pres

def add_plot_slide(pres, conf, plot_type):
    """
    Creates matplotlib or pptx plot slide in pres based on conf

    Prameters: 
    pres (Presentation): pptx Presentation object
    conf (dict): parameters for slide creation
    plot_type (str): type of plot (pptx or matplotlib)
    """
    title_slide_layout = pres.slide_layouts[5]
    slide = pres.slides.add_slide(title_slide_layout)
    t = slide.shapes.title
    t.text = conf["title"]
    df = pd.read_csv(conf["content"], header = None, delimiter = ";")
    df.columns = ["x", "y"]
    if plot_type == "pptx":
        y = df["y"].to_list()
        chart_data = ChartData()
        chart_data.categories = df["x"].to_list()
        chart_data.add_series("Series", y)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(2), Inches(2), Inches(6), Inches(4.5), chart_data).chart
        chart.has_legend = False
        
        category_axis_title = chart.category_axis.axis_title

        category_axis_title.text_frame.text = conf["configuration"]["x-label"]
        value_axis_title = chart.value_axis.axis_title
        value_axis_title.text_frame.text = conf["configuration"]["y-label"]
    if plot_type == "matplotlib":
        plt.plot(df.x, df.y)
        image_stream = io.BytesIO()
        plt.savefig(image_stream)
        pic = slide.shapes.add_picture(image_stream, Inches(0.5), Inches(1.75), width=Inches(9), height=Inches(5))


def check_input(input):
    if input != "pptx" and input != "matplotlib":
        return False
    return True

def main():
    print("Input file name (json):")
    json_input = input()
    data = read_config(json_input)
    print("Plot type (pptx or matplotlib):")
    plot_type = input()
    if not check_input(plot_type):
        print("Wrong plot type")
        return

    prs = Presentation()
    


    for conf in data["presentation"]:
        if conf["type"] == "title":
            add_title_slide(prs, conf)
        if conf["type"] == "text":
            add_text_slide(prs, conf)        
        if conf["type"] == "list":
            add_list_slide(prs, conf)
        if conf["type"] == "picture":
            add_picture_slide(prs, conf)
        if conf["type"] == "plot":
            add_plot_slide(prs, conf, plot_type)


    prs.save('presentation.pptx')

if __name__ == "__main__":
    main()
